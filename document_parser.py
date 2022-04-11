import pickle
import re
import os
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

from modules.models.segmentor import cleaner

from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LAParams, LTTextBox
from pdfminer.pdfinterp import PDFPageInterpreter, PDFResourceManager
from pdfminer.pdfpage import PDFPage

done_path = "data/documentos/done"


def process_pdfs(path: str):
    """
    Function to parse pdf documents
    First, check if there is a local db stored.
    Then, given a folder path, get the files inside it and do the following for each pdf file:
        - Get all data information that is outside tables and images.
        - Process all strings and deletes errors from the parser.
        - Divide the data into paragraphs and store the information in a Dict type Object.
    At the end, store all documents information into a pickle object to use it on the next step.
    """

    # Stored DataBase check
    boolie = False
    if len(os.listdir("data/pickle")) > 1:
        with open("data/pickle/document_list.pkl", "rb") as f:
            data = pickle.load(f)
            boolie = True

    # Creation of the document wrapper
    document_list = []

    folder = os.listdir(path)
    asignatura = path[16:].title()

    # PDF files format check
    for pdf in folder:
        if pdf[-4:] != ".pdf":
            continue
        print(pdf)

        # PDFMiner settings
        fp = open(f"{path}/{pdf}", "rb")
        rsrcmgr = PDFResourceManager()
        laparams = LAParams()
        device = PDFPageAggregator(rsrcmgr, laparams=laparams)
        interpreter = PDFPageInterpreter(rsrcmgr, device)
        pages = PDFPage.get_pages(fp)

        # Information parsing
        tuplas = {}
        resultado = []
        for i, page in enumerate(pages):

            if i == 1:
                continue

            tuplas[i] = []
            interpreter.process_page(page)
            layout = device.get_result()
            for lobj in layout:
                if isinstance(lobj, LTTextBox):
                    x, y, text = lobj.bbox[0], lobj.bbox[3], lobj.get_text()

                    if y > 70:
                        text = re.sub("_{2,}.+(?:\n+.+)+", "", text)
                        text = text.replace("(cid:4)", "")
                        tuplas[i].append((y, text))

            prueba = []
            for a in tuplas[i]:
                prueba.append(a)
            prueba.sort(key=lambda y: (-y[0]))

            pagina = []
            for a in prueba:

                pagina.append(" ".join(a[1].split()))

            resultado.append(pagina)
            no_blanks = []
        for pagina in resultado:
            nano = []
            for parrafo in pagina:
                if len(parrafo) > 1:
                    clean = re.sub("•", "", parrafo).strip()
                    nano.append(clean)
            no_blanks.append(nano)

        for a in no_blanks:
            texto = " ".join(a)
            document_list.append({"title": f"{asignatura} - {pdf}", "content": texto})

        # Moving pdf parsed to done:
        fp.close()
        os.replace(f"{path}/{pdf}", f"{done_path}/{pdf}")

    if boolie:
        data.extend(document_list)
    else:
        data = document_list

    # Create the pkl to use in next step
    with open("data/pickle/document_list.pkl", "wb") as f:
        pickle.dump(data, f)


def process_pptx(path: str):
    """
    Function to parse pptx documents.
    First, check if there is a local db stored.
    Then, given a folder path, get the files inside it and do the following for each pptx file:
        - Get all data information that is outside images.
        - Process all strings and deletes encoding-based errors from the parser.
        - Divide the data into slides and store the information in a Dict type Object.
    At the end, store all documents information into a pickle object to use it on the next step.
    """

    asignatura = path[16:]
    salida = []
    folder = os.listdir(path)
    boolie = False
    if len(os.listdir("data/pickle")) > 1:
        with open("data/pickle/document_list.pkl", "rb") as f:
            data = pickle.load(f)
            boolie = True

    for eachfile in folder:
        if eachfile[-5:] != ".pptx":
            continue
        print(eachfile)

        with open(f"{path}/{eachfile}", "rb") as f:
            prs = Presentation(f)

        datos_powerpoint = []
        for i, slide in enumerate(prs.slides):
            datos_slide = []

            # Jump to slide nº 3
            if i not in [0, 1]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):

                        datos_slide.append(shape.text)
            datos_powerpoint.append(datos_slide)

        for a in datos_powerpoint:
            if len(a) > 0:
                text = " ".join(a)
                text = text.replace("\t", " ")
                text = text.replace("\n", " ")
                text = text.replace("\x0b", " ")
                text = text.replace("\uf075", " ")
                text = text.replace("\xa0", " ")
                text = " ".join(text.split()).strip()
                if len(text.split()) > 3:
                    salida.append(
                        {"title": f"{asignatura} - {eachfile}", "content": text}
                    )
        os.replace(f"{path}/{eachfile}", f"{done_path}/{eachfile}")
    if boolie:
        data.extend(salida)
    else:
        data = salida
    # Create the pkl to use in next step
    with open("data/pickle/document_list.pkl", "wb") as f:
        pickle.dump(data, f)


def process_html(url: str):
    """
    Function to parse html documents
    The function check if there is a local db stored.
    Then, given a URL, try to scrap the web and access its information:
        - Focus on the body tag.
        - Divide the data into paragraphs using <h> tags.
        - Deletes all the other tags inside it and store the information in a Dict type Object.
    At the end, store all documents information into a pickle object to use it on the next step.
    """
    tema = re.split("\/", url)[-1]

    r = requests.get(url)
    soup = BeautifulSoup(r.text, "html.parser")
    body = str(soup.body)

    salida = []
    pa = re.split("<h\d>", body)
    for i in pa:
        if len(cleaner(i)) > 1:
            salida.append({"title": tema, "content": (" ".join(cleaner(i)))})

    if len(os.listdir("data/pickle")) > 1:
        with open("data/pickle/document_list.pkl", "rb") as f:
            data = pickle.load(f)
            salida.extend(data)

    with open("data/pickle/document_list.pkl", "wb") as f:
        pickle.dump(salida, f)

    with open(f"{done_path}/{tema}.txt", "a") as f:
        f.write(url)


def process_path(path: str):
    # Parent function that call its childs

    if path[0] != "d":
        process_html(path)
    else:
        process_pdfs(path)
        process_pptx(path)


process_path("https://ariesco.github.io/OIM/docs/intro.html")
